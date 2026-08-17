"""SMB/CIFS file-share content scanner (Surface add-on).

Crawls a Windows file share, extracts the *body* of documents
(PDF/Word/Excel/PPTX + text/config/code) and flags confidential data:
a built-in high-signal secret ruleset plus per-target custom regex.

Design notes:
- Credentials: per-target (asset.config.smb_username/smb_domain + an
  AES-encrypted smb_password_enc, decrypted at scan time) take precedence
  over the SURFACE_SMB_USERNAME / SURFACE_SMB_PASSWORD / SURFACE_SMB_DOMAIN
  service-account fallback. The plaintext password is never stored.
- Findings keep `target` = the file's UNC path (so the file is visible and
  dedup stays per-file); `type` = the matched rule name, so several distinct
  secrets in one file are distinct findings. The UI rolls findings up to the
  host matching the share's server name.
- Incremental: a SQLite cache (stdlib, no migration) keyed by
  (path, size, mtime) skips files unchanged since the last scan — only the
  first scan reads everything. Cache dir = SURFACE_SMB_CACHE_DIR (default
  /tmp/smb-cache; mount a volume there for persistence across restarts).
- Findings never store the raw secret: the match is masked.
- Registered with the core via the SURFACE_SCANNERS contract (wants_config),
  so it gets (value, asset_config).
"""
from __future__ import annotations

import hashlib
import logging
import os
import re
import sqlite3
import time
from concurrent.futures import ThreadPoolExecutor, as_completed

logger = logging.getLogger("surface-smb")

# ── Defaults (overridable per target via asset.config) ────────────
_DEFAULT_EXTENSIONS = {
    "pdf", "doc", "docx", "xls", "xlsx", "ppt", "pptx",
    "txt", "csv", "log", "ini", "conf", "cfg", "config", "env",
    "yml", "yaml", "json", "xml", "properties", "sh", "ps1", "bat",
    "sql", "key", "pem", "kdbx", "rdp",
}
_DEFAULT_MAX_SIZE_MB = 50
# Workers default kept LOW: document extraction (pdfminer/openpyxl/python-pptx)
# is CPU-bound and holds the GIL. The whole scan runs in a worker thread of the
# FastAPI process, so too many extractor threads starve the asyncio event loop
# and make the whole app lag for the duration of a big scan. 3 is a safe floor.
_MAX_WORKERS = int(os.getenv("SURFACE_SMB_WORKERS", "3") or "3")
# Hard safety caps so a large filer can't run for hours / hog the box. Tunable
# per target (config.max_files / config.time_budget_s) or via env.
_DEFAULT_MAX_FILES = int(os.getenv("SURFACE_SMB_MAX_FILES", "5000") or "5000")
_DEFAULT_TIME_BUDGET_S = int(os.getenv("SURFACE_SMB_TIME_BUDGET_S", "1800") or "1800")
# Small pause between extraction batches to yield the GIL back to the event
# loop, so the app stays responsive while a scan runs. Seconds.
_BATCH_PAUSE_S = float(os.getenv("SURFACE_SMB_BATCH_PAUSE_S", "0.05") or "0.05")

# ── Built-in secret ruleset (name, severity, compiled regex) ──────
# High-signal patterns; kept deliberately tight to limit false positives.
_SECRET_RULES: list[tuple[str, str, "re.Pattern[str]"]] = [
    ("private_key", "critical", re.compile(r"-----BEGIN (?:RSA |EC |OPENSSH |DSA |PGP )?PRIVATE KEY-----")),
    ("aws_access_key", "critical", re.compile(r"\b(?:AKIA|ASIA)[0-9A-Z]{16}\b")),
    ("aws_secret_key", "high", re.compile(r"\baws_secret_access_key\s*[=:]\s*['\"]?[A-Za-z0-9/+]{40}\b")),
    ("gcp_service_account", "critical", re.compile(r'"type"\s*:\s*"service_account"')),
    ("google_api_key", "high", re.compile(r"\bAIza[0-9A-Za-z_\-]{35}\b")),
    ("github_token", "critical", re.compile(r"\bgh[pousr]_[A-Za-z0-9]{36,}\b")),
    ("slack_token", "high", re.compile(r"\bxox[baprs]-[A-Za-z0-9-]{10,}\b")),
    ("jwt", "medium", re.compile(r"\beyJ[A-Za-z0-9_\-]{10,}\.[A-Za-z0-9_\-]{10,}\.[A-Za-z0-9_\-]{10,}\b")),
    ("connection_string", "high", re.compile(r"(?i)(?:Data Source|Server)=[^;]+;[^;]*(?:Password|Pwd)=[^;]+")),
    ("password_assignment", "medium", re.compile(r"(?i)(?:password|passwd|pwd|mot_de_passe|motdepasse)\s*[=:]\s*['\"]?[^\s'\"]{6,}")),
    ("generic_secret", "medium", re.compile(r"(?i)\b(?:secret|api[_-]?key|token|client[_-]?secret)\s*[=:]\s*['\"]?[A-Za-z0-9/+_\-]{16,}")),
]

# Snaffler-style interesting filenames (flagged even without a content hit).
_INTERESTING_NAME = re.compile(
    r"(?i)(?:^|[\\/])(?:id_rsa|id_dsa|id_ed25519|\.npmrc|\.pgpass|\.htpasswd|"
    r"web\.config|unattend\.xml|sysprep\.xml|wp-config\.php|\.kdbx|"
    r"credentials|secrets?|password|motdepasse|backup)\b"
)

_MAX_FINDINGS_PER_FILE = 20


def _mask(s: str) -> str:
    s = s.strip()
    if len(s) <= 8:
        return s[0] + "***" if s else "***"
    return f"{s[:4]}…{s[-2:]} ({len(s)} chars)"


# ── Share path parsing ────────────────────────────────────────────
def _parse_share(value: str) -> tuple[str, str]:
    """'\\\\host\\share\\sub' | '//host/share' | 'smb://host/share' ->
    (server, unc_root) where unc_root is a back-slash UNC path."""
    raw = (value or "").strip()
    low = raw.lower()
    if low.startswith("smb://"):
        raw = raw[6:]
    norm = raw.replace("\\", "/").lstrip("/")
    parts = [p for p in norm.split("/") if p]
    if len(parts) < 2:
        raise ValueError(f"Invalid share path: {value!r} (expected \\\\host\\share)")
    server = parts[0]
    unc_root = "\\\\" + "\\".join(parts)
    return server, unc_root


# ── Incremental cache (SQLite) ────────────────────────────────────
class _Cache:
    def __init__(self, share_key: str):
        cache_dir = os.getenv("SURFACE_SMB_CACHE_DIR", "/tmp/smb-cache")
        try:
            os.makedirs(cache_dir, exist_ok=True)
            digest = hashlib.sha256(share_key.encode()).hexdigest()[:16]
            self.conn: sqlite3.Connection | None = sqlite3.connect(os.path.join(cache_dir, f"{digest}.db"))
            self.conn.execute(
                "CREATE TABLE IF NOT EXISTS scanned (path TEXT PRIMARY KEY, size INTEGER, mtime REAL)"
            )
            self.conn.commit()
        except Exception as e:  # cache is best-effort — degrade to full scan
            logger.warning("SMB cache unavailable (%s): scanning all files", e)
            self.conn = None

    def unchanged(self, path: str, size: int, mtime: float) -> bool:
        if not self.conn:
            return False
        row = self.conn.execute("SELECT size, mtime FROM scanned WHERE path=?", (path,)).fetchone()
        return bool(row and row[0] == size and abs(row[1] - mtime) < 1.0)

    def mark(self, path: str, size: int, mtime: float) -> None:
        if not self.conn:
            return
        self.conn.execute("INSERT OR REPLACE INTO scanned VALUES (?,?,?)", (path, size, mtime))

    def commit(self) -> None:
        if self.conn:
            try:
                self.conn.commit()
                self.conn.close()
            except Exception:
                pass


# ── Document text extraction ──────────────────────────────────────
def _extract_text(ext: str, data: bytes) -> str:
    import io
    try:
        if ext == "pdf":
            from pdfminer.high_level import extract_text
            return extract_text(io.BytesIO(data)) or ""
        if ext == "docx":
            import docx
            d = docx.Document(io.BytesIO(data))
            return "\n".join(p.text for p in d.paragraphs)
        if ext == "xlsx":
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
            out = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    out.append(" ".join(str(c) for c in row if c is not None))
            wb.close()
            return "\n".join(out)
        if ext == "pptx":
            from pptx import Presentation
            prs = Presentation(io.BytesIO(data))
            out = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        out.append(shape.text_frame.text)
            return "\n".join(out)
        # plain text / config / code
        return data.decode("utf-8", errors="replace")
    except Exception as e:
        logger.info("extract failed (%s): %s", ext, e)
        return ""


# ── Content scan ──────────────────────────────────────────────────
# Chunk size for operator-supplied regexes — see _scan_text.
_CUSTOM_RE_CHUNK = 8192
_CUSTOM_RE_OVERLAP = 512

def _scan_text(text: str, custom: list["re.Pattern[str]"]) -> list[tuple[str, str, str]]:
    """Return [(rule_name, severity, masked_match)] capped per file."""
    hits: list[tuple[str, str, str]] = []
    for name, sev, rx in _SECRET_RULES:
        for m in rx.finditer(text):
            hits.append((name, sev, _mask(m.group(0))))
            if len(hits) >= _MAX_FINDINGS_PER_FILE:
                return hits
    # Operator-supplied patterns run over file contents, so a pattern with
    # nested quantifiers ("(a+)+$") backtracks catastrophically and pins a
    # worker. CPython cannot interrupt a match once started — there is no
    # timeout on `re` — so the only dependency-free lever is the SIZE of each
    # match attempt: backtracking blows up with input length, and chunking
    # caps how long any single attempt can run. Overlap keeps a match that
    # straddles a boundary from being missed.
    #
    # This bounds the damage, it does not eliminate the class. A robust fix
    # needs a linear-time engine (`re2`) or `regex` with `timeout=` — both new
    # dependencies. Note the setting is admin-only since the Surface routes
    # were role-gated, so the realistic threat is a misconfiguration, not an
    # anonymous denial of service.
    for rx in custom:
        for start in range(0, len(text) or 1, _CUSTOM_RE_CHUNK):
            chunk = text[start:start + _CUSTOM_RE_CHUNK + _CUSTOM_RE_OVERLAP]
            for m in rx.finditer(chunk):
                hits.append(("custom_regex", "medium", _mask(m.group(0))))
                if len(hits) >= _MAX_FINDINGS_PER_FILE:
                    return hits
    return hits


def _finding(severity: str, title: str, description: str, target: str, evidence: dict,
             type_: str = "sensitive_data") -> dict:
    # `type_` participates in the dedup key (scanner|type|target). Passing the
    # matched rule name makes each (file, rule) pair a distinct logical finding,
    # so several secrets in the same file don't collapse into one row.
    return {"scanner": "smb_scan", "type": type_, "severity": severity,
            "title": title, "description": description, "target": target, "evidence": evidence}


def scan_smb_share(value: str, config: dict | None = None) -> list[dict]:
    config = config or {}
    try:
        import smbclient
    except ImportError:
        return [_finding("info", "smbprotocol not installed", "The SMB add-on requires smbprotocol.", value, {}, type_="smb_status")]

    try:
        server, unc_root = _parse_share(value)
    except ValueError as e:
        return [_finding("info", f"Invalid share: {value}", str(e), value, {}, type_="smb_status")]

    # Per-target credentials (from the UI, password decrypted) take precedence
    # over the SURFACE_SMB_* default service account.
    user = (config.get("smb_username") or os.getenv("SURFACE_SMB_USERNAME", "")).strip()
    domain = (config.get("smb_domain") or os.getenv("SURFACE_SMB_DOMAIN", "")).strip()
    pw = os.getenv("SURFACE_SMB_PASSWORD", "")
    if config.get("smb_password_enc"):
        try:
            from src.crypto import decrypt_secret
            pw = decrypt_secret(config["smb_password_enc"]) or pw
        except Exception as e:
            logger.warning("SMB password decrypt failed: %s", e)
    if not user or not pw:
        return [_finding("info", "SMB credentials missing",
                         "Set the credentials on the target, or define "
                         "SURFACE_SMB_USERNAME / SURFACE_SMB_PASSWORD (+ DOMAIN).", value, {},
                         type_="smb_status")]
    full_user = f"{domain}\\{user}" if domain else user

    exts = {e.lower().lstrip(".") for e in config.get("extensions") or []} or _DEFAULT_EXTENSIONS
    max_bytes = int(config.get("max_size_mb") or _DEFAULT_MAX_SIZE_MB) * 1024 * 1024
    max_files = int(config.get("max_files") or _DEFAULT_MAX_FILES)
    time_budget = int(config.get("time_budget_s") or _DEFAULT_TIME_BUDGET_S)
    workers = max(1, int(config.get("workers") or _MAX_WORKERS))
    deadline = time.monotonic() + time_budget
    truncated: str | None = None
    custom = []
    for pat in config.get("custom_regex") or []:
        # Reject the textbook catastrophic shapes outright — a quantifier
        # applied to an already-quantified group. Heuristic, not a proof: it
        # catches the patterns people actually write by accident.
        if re.search(r"\([^)]*[+*][^)]*\)\s*[+*]", pat or ""):
            logger.warning("refused nested-quantifier regex %r (ReDoS risk)", pat)
            continue
        try:
            custom.append(re.compile(pat))
        except re.error as e:
            logger.warning("bad custom regex %r: %s", pat, e)

    findings: list[dict] = []
    cache = _Cache(unc_root)
    try:
        smbclient.register_session(server, username=full_user, password=pw)
    except Exception as e:
        cache.commit()
        return [_finding("info", f"SMB connection failed: {server}", type(e).__name__, value, {}, type_="smb_status")]

    # 1. Enumerate candidate files (path, size, mtime), skipping unchanged.
    #    Bounded by max_files and the time budget so a large filer can't run
    #    for hours.
    candidates: list[tuple[str, int]] = []
    scanned = skipped = 0
    try:
        for dirpath, _dirs, files in smbclient.walk(unc_root):
            if time.monotonic() > deadline:
                truncated = "time"
                break
            for fn in files:
                if len(candidates) >= max_files:
                    truncated = truncated or "files"
                    break
                ext = fn.rsplit(".", 1)[-1].lower() if "." in fn else ""
                fpath = dirpath.rstrip("\\") + "\\" + fn
                if ext not in exts and not _INTERESTING_NAME.search(fpath):
                    continue
                try:
                    st = smbclient.stat(fpath)
                except Exception:
                    continue
                if st.st_size > max_bytes:
                    continue
                if cache.unchanged(fpath, st.st_size, st.st_mtime):
                    skipped += 1
                    continue
                # interesting filename → flag regardless of content
                is_name = bool(_INTERESTING_NAME.search(fpath))
                if is_name:
                    findings.append(_finding("low", f"Sensitive file by name: {fn}",
                                             "Name/extension suggestive of sensitive data.", fpath,
                                             {"file": fpath, "rule": "interesting_name"},
                                             type_="interesting_name"))
                if ext in exts:
                    # Content files are cached only AFTER they are read+extracted
                    # (below), so a truncated scan resumes the unread ones.
                    candidates.append((fpath, st.st_size, st.st_mtime))
                elif is_name:
                    # Name-only hit: fully handled now, safe to cache.
                    cache.mark(fpath, st.st_size, st.st_mtime)
            if truncated == "files":
                break
    except Exception as e:
        logger.exception("SMB walk failed for %s", unc_root)
        findings.append(_finding("info", f"Share crawl interrupted: {server}", type(e).__name__, value, {}, type_="smb_status"))

    # 2. Read + extract + match, in parallel (I/O-bound).
    def _process(item: tuple[str, int, float]) -> list[dict]:
        fpath = item[0]
        ext = fpath.rsplit(".", 1)[-1].lower() if "." in fpath else ""
        try:
            with smbclient.open_file(fpath, mode="rb") as fh:
                data = fh.read()
        except Exception:
            return []
        text = _extract_text(ext, data)
        if not text:
            return []
        out = []
        for name, sev, masked in _scan_text(text, custom):
            out.append(_finding(sev, f"Sensitive data ({name}): {fpath.rsplit(chr(92), 1)[-1]}",
                                f"Pattern \"{name}\" detected in the file body.", fpath,
                                {"file": fpath, "rule": name, "match": masked},
                                type_=name))
        return out

    # 2b. Read + extract + match in bounded batches. Between batches we check
    #     the time budget and briefly yield the GIL so the FastAPI event loop
    #     stays responsive during a long scan.
    if candidates:
        batch = max(1, workers) * 4
        with ThreadPoolExecutor(max_workers=workers) as pool:
            for i in range(0, len(candidates), batch):
                if time.monotonic() > deadline:
                    truncated = truncated or "time"
                    break
                chunk = candidates[i:i + batch]
                futs = {pool.submit(_process, c): c for c in chunk}
                for fut in as_completed(futs):
                    item = futs[fut]
                    try:
                        findings.extend(fut.result())
                        scanned += 1
                        # Mark as scanned only now (this version is fully handled).
                        # Done in the main thread — the sqlite cache is not
                        # safe to touch from the pool's worker threads.
                        cache.mark(item[0], item[1], item[2])
                    except Exception:
                        pass
                if _BATCH_PAUSE_S > 0:
                    time.sleep(_BATCH_PAUSE_S)

    cache.commit()
    try:
        smbclient.delete_session(server)
    except Exception:
        pass

    if truncated:
        reason = "time budget reached" if truncated == "time" else f"cap of {max_files} files reached"
        findings.append(_finding(
            "info", f"Share scan truncated: {server}",
            f"Scan interrupted ({reason}). {scanned} file(s) read, {skipped} unchanged skipped. "
            "The next scan resumes the remaining files (incremental cache). "
            "To cover more in one pass: increase max_files / time_budget_s, "
            "or refine the extensions / max size.",
            value, {"scanned": scanned, "skipped": skipped, "limit": truncated, "max_files": max_files},
            type_="smb_status"))

    logger.info("SMB scan %s: %d files read, %d skipped (cache), %d findings, truncated=%s",
                unc_root, scanned, skipped, len(findings), truncated)
    return findings


# ── In-app help documentation (served by the core only when this add-on is
# loaded — see scanners.addon_help_docs + GET /api/help/addons). The HTML is
# injected into the Méthodologie / Utilisation help tabs. Because the text
# lives HERE (not in the core i18n bundle), an image built without this add-on
# never ships it. Keep it developer-authored HTML only (no user input).
_DOC = {
    "fr": {
        "methodo":
            '<h2>Scan de partages de fichiers (SMB/CIFS) <em>— add-on</em></h2>'
            '<p>Cet add-on étend Surface au-delà de la surface externe : il inspecte le <strong>contenu</strong> des partages Windows (SMB/CIFS) à la recherche de données confidentielles exposées en interne (secrets, identifiants, fichiers sensibles).</p>'
            '<ul>'
            '<li><strong>Parcours du partage</strong> via le protocole SMB avec un compte de service en lecture seule, filtrage par extension et taille maximale par fichier.</li>'
            '<li><strong>Extraction du corps des documents</strong> — le texte est extrait du contenu réel des <strong>PDF, Word (.docx), Excel (.xlsx), PowerPoint (.pptx)</strong>, en plus des fichiers texte / config / code. Un secret enfoui dans un document bureautique est donc détecté, pas seulement dans les fichiers plats.</li>'
            '<li><strong>Ruleset de secrets intégré</strong> — clés privées (RSA/EC/OpenSSH), clés AWS/GCP/Google, tokens GitHub/Slack, JWT, chaînes de connexion base de données, affectations de mot de passe, secrets génériques (<code>api_key</code>, <code>token</code>, <code>client_secret</code>…). Chaque correspondance est <strong>masquée</strong> : le secret n\'est jamais stocké.</li>'
            '<li><strong>Regex personnalisées par cible</strong> — en complément du ruleset, on ajoute ses propres motifs (numéros de dossier, identifiants métier…).</li>'
            '<li><strong>Noms de fichiers sensibles</strong> (style Snaffler) — <code>id_rsa</code>, <code>.kdbx</code>, <code>web.config</code>, <code>*credentials*</code>, <code>*password*</code>, <code>*backup*</code>… sont flaggés même sans correspondance de contenu.</li>'
            '<li><strong>Scan incrémental</strong> — un cache (chemin + taille + date de modification) évite de relire les fichiers inchangés : seul le premier scan lit tout, les suivants ne traitent que les nouveautés.</li>'
            '</ul>'
            '<div class="help-tip"><strong>Rattachement :</strong> les findings d\'un partage <code>\\\\serveur\\partage</code> remontent sur le host <code>serveur</code> dans la vue <strong>Hosts</strong>, comme les autres scans. Le <code>type</code> de chaque finding porte la règle déclenchée, donc plusieurs secrets distincts d\'un même fichier restent des findings distincts.</div>'
            '<div class="help-tip"><strong>Confidentialité :</strong> scan en lecture seule, le contenu des fichiers n\'est jamais conservé — seuls le chemin, la règle déclenchée et un extrait <em>masqué</em> du secret sont stockés.</div>',
        "usage":
            '<h2>Cible « Partage de fichiers » <em>(add-on SMB)</em></h2>'
            '<p>Quand l\'add-on SMB est installé, un type de cible supplémentaire apparaît à l\'ajout : <strong>Partage de fichiers</strong>.</p>'
            '<ul>'
            '<li><strong>Chemin du partage</strong> — au format <code>\\\\serveur\\partage</code> (ou <code>//serveur/partage</code>, <code>smb://serveur/partage</code>). Un sous-dossier peut être précisé.</li>'
            '<li><strong>Identifiants</strong> — identifiant SMB, domaine (vide pour un compte local) et mot de passe d\'un compte de service en <strong>lecture seule</strong>. Le mot de passe est <strong>chiffré</strong> côté serveur et jamais renvoyé par l\'API ; lors d\'une modification, laissez le champ vide pour le conserver.</li>'
            '<li><strong>Options</strong> — extensions à scanner, taille maximale par fichier, et regex personnalisées (une par ligne).</li>'
            '</ul>'
            '<p>Après le scan, le partage apparaît dans la vue <strong>Hosts</strong> sous le nom de son serveur (badge « partage »), avec ses findings regroupés comme pour n\'importe quel host.</p>',
    },
    "en": {
        "methodo":
            '<h2>File-share scanning (SMB/CIFS) <em>— add-on</em></h2>'
            '<p>This add-on extends Surface beyond the external attack surface: it inspects the <strong>contents</strong> of Windows file shares (SMB/CIFS) for confidential data exposed internally (secrets, credentials, sensitive files).</p>'
            '<ul>'
            '<li><strong>Share crawl</strong> over the SMB protocol with a read-only service account, filtered by extension and per-file size cap.</li>'
            '<li><strong>Document body extraction</strong> — text is pulled from the actual content of <strong>PDF, Word (.docx), Excel (.xlsx), PowerPoint (.pptx)</strong>, on top of text / config / code files. A secret buried inside an office document is detected, not just in flat files.</li>'
            '<li><strong>Built-in secret ruleset</strong> — private keys (RSA/EC/OpenSSH), AWS/GCP/Google keys, GitHub/Slack tokens, JWTs, database connection strings, password assignments, generic secrets (<code>api_key</code>, <code>token</code>, <code>client_secret</code>…). Every match is <strong>masked</strong>: the secret is never stored.</li>'
            '<li><strong>Per-target custom regex</strong> — on top of the ruleset, add your own patterns (case numbers, business identifiers…).</li>'
            '<li><strong>Sensitive file names</strong> (Snaffler-style) — <code>id_rsa</code>, <code>.kdbx</code>, <code>web.config</code>, <code>*credentials*</code>, <code>*password*</code>, <code>*backup*</code>… are flagged even without a content match.</li>'
            '<li><strong>Incremental scan</strong> — a cache (path + size + mtime) skips unchanged files: only the first scan reads everything, later scans process only what changed.</li>'
            '</ul>'
            '<div class="help-tip"><strong>Roll-up:</strong> findings from a share <code>\\\\server\\share</code> surface under the host <code>server</code> in the <strong>Hosts</strong> view, like every other scan. Each finding\'s <code>type</code> carries the rule it triggered, so several distinct secrets in one file stay distinct findings.</div>'
            '<div class="help-tip"><strong>Privacy:</strong> read-only scan, file contents are never kept — only the path, the triggered rule and a <em>masked</em> snippet of the secret are stored.</div>',
        "usage":
            '<h2>"File share" target <em>(SMB add-on)</em></h2>'
            '<p>When the SMB add-on is installed, an extra target type appears when adding an asset: <strong>File share</strong>.</p>'
            '<ul>'
            '<li><strong>Share path</strong> — as <code>\\\\server\\share</code> (or <code>//server/share</code>, <code>smb://server/share</code>). A sub-folder can be specified.</li>'
            '<li><strong>Credentials</strong> — SMB username, domain (empty for a local account) and the password of a <strong>read-only</strong> service account. The password is <strong>encrypted</strong> server-side and never returned by the API; when editing, leave the field empty to keep it.</li>'
            '<li><strong>Options</strong> — extensions to scan, max size per file, and custom regex (one per line).</li>'
            '</ul>'
            '<p>After scanning, the share appears in the <strong>Hosts</strong> view under its server name (badge "share"), with its findings grouped like any other host.</p>',
    },
}

SURFACE_SCANNERS = {
    "smb_scan": {
        "label": "SMB file-share content (secrets & sensitive data)",
        "kinds": {"file_share"},
        "callable": scan_smb_share,
        "returns_discovered": False,
        "wants_config": True,
        "doc": _DOC,
    },
}
SURFACE_DEFAULT_SCANNERS = {"file_share": ["smb_scan"]}
